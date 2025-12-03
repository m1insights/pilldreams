"""
Exports API Endpoints
Week 5: Excel and PowerPoint export functionality
Enhanced with competitive landscape PPTX exports
"""

from fastapi import APIRouter, HTTPException, Depends, Request
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional, List, Literal
from datetime import datetime
import os
import io
from supabase import create_client, Client
from dotenv import load_dotenv

load_dotenv()

router = APIRouter(prefix="/exports", tags=["exports"])

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY") or os.getenv("SUPABASE_KEY")
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)


# ============================================
# Pydantic Models
# ============================================

class ExportRequest(BaseModel):
    entity_type: str  # 'drugs', 'targets', 'trials', 'scores', 'watchlist'
    entity_ids: Optional[List[str]] = None  # Specific IDs, or None for all
    format: str = "xlsx"  # 'xlsx' or 'csv'
    include_scores: bool = True
    include_chemistry: bool = False
    include_trials: bool = False


class DealMemoRequest(BaseModel):
    drug_id: str
    indication_id: Optional[str] = None
    include_competition: bool = True
    include_trials: bool = True
    include_chemistry: bool = True
    template: str = "standard"  # 'standard', 'executive', 'detailed'


class LandscapeExportRequest(BaseModel):
    """Request for competitive landscape PowerPoint export."""
    export_type: Literal["target", "indication", "company", "pipeline"] = "target"
    target_id: Optional[str] = None  # For target landscape
    indication_id: Optional[str] = None  # For indication landscape
    company_id: Optional[str] = None  # For company portfolio
    drug_ids: Optional[List[str]] = None  # For custom pipeline selection
    include_scores: bool = True
    include_trials: bool = True
    include_chemistry: bool = False
    template: Literal["executive", "detailed", "comparison"] = "executive"


# ============================================
# Helper: Auth
# ============================================

async def get_current_user(request: Request) -> Optional[str]:
    """Extract user from Authorization header."""
    auth_header = request.headers.get("Authorization")
    if not auth_header or not auth_header.startswith("Bearer "):
        return None
    token = auth_header.replace("Bearer ", "")
    try:
        user = supabase.auth.get_user(token)
        if user and user.user:
            return user.user.id
    except Exception:
        pass
    return None


async def require_auth(request: Request) -> str:
    """Require authentication."""
    user_id = await get_current_user(request)
    if not user_id:
        raise HTTPException(status_code=401, detail="Not authenticated")
    return user_id


async def check_export_access(user_id: str) -> bool:
    """Check if user can export (has exports remaining)."""
    try:
        profile = supabase.table("ci_user_profiles").select(
            "exports_this_month, exports_limit, subscription_tier"
        ).eq("id", user_id).execute()

        if not profile.data:
            return False

        user = profile.data[0]
        limit = user.get("exports_limit", 5)

        # -1 means unlimited
        if limit == -1:
            return True

        return user.get("exports_this_month", 0) < limit
    except Exception:
        return True  # Allow if can't check (table doesn't exist)


async def increment_export_count(user_id: str):
    """Increment user's export count."""
    try:
        supabase.rpc("increment_usage", {
            "p_user_id": user_id,
            "p_counter": "exports"
        }).execute()
    except Exception:
        # Fallback: direct update
        try:
            supabase.table("ci_user_profiles").update({
                "exports_this_month": supabase.table("ci_user_profiles")
                    .select("exports_this_month")
                    .eq("id", user_id)
                    .execute()
                    .data[0]["exports_this_month"] + 1
            }).eq("id", user_id).execute()
        except Exception:
            pass


# ============================================
# Excel Export
# ============================================

@router.post("/excel")
async def export_to_excel(
    data: ExportRequest,
    request: Request,
    user_id: str = Depends(require_auth)
):
    """Export data to Excel format."""
    # Check export access
    if not await check_export_access(user_id):
        raise HTTPException(
            status_code=403,
            detail="Export limit reached. Upgrade to Pro for more exports."
        )

    try:
        import pandas as pd
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils.dataframe import dataframe_to_rows
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="Excel export requires openpyxl. Run: pip install openpyxl"
        )

    # Fetch data based on entity type
    if data.entity_type == "drugs":
        df = await fetch_drugs_data(data.entity_ids, data.include_scores)
        filename = "phase4_drugs_export"
    elif data.entity_type == "targets":
        df = await fetch_targets_data(data.entity_ids)
        filename = "phase4_targets_export"
    elif data.entity_type == "trials":
        df = await fetch_trials_data(data.entity_ids)
        filename = "phase4_trials_export"
    elif data.entity_type == "scores":
        df = await fetch_scores_data(data.entity_ids)
        filename = "phase4_scores_export"
    elif data.entity_type == "watchlist":
        df = await fetch_watchlist_data(user_id)
        filename = "phase4_watchlist_export"
    else:
        raise HTTPException(status_code=400, detail=f"Invalid entity type: {data.entity_type}")

    if df.empty:
        raise HTTPException(status_code=404, detail="No data found for export")

    # Create Excel workbook with styling
    wb = Workbook()
    ws = wb.active
    ws.title = data.entity_type.title()

    # Header style
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="1a1a1a", end_color="1a1a1a", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center")

    # Add data
    for r_idx, row in enumerate(dataframe_to_rows(df, index=False, header=True), 1):
        for c_idx, value in enumerate(row, 1):
            cell = ws.cell(row=r_idx, column=c_idx, value=value)
            if r_idx == 1:  # Header row
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_alignment

    # Auto-adjust column widths
    for column in ws.columns:
        max_length = 0
        column_letter = column[0].column_letter
        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width

    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)

    # Increment export count
    await increment_export_count(user_id)

    # Return file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={
            "Content-Disposition": f"attachment; filename={filename}_{timestamp}.xlsx"
        }
    )


@router.post("/csv")
async def export_to_csv(
    data: ExportRequest,
    request: Request,
    user_id: str = Depends(require_auth)
):
    """Export data to CSV format."""
    # Check export access
    if not await check_export_access(user_id):
        raise HTTPException(
            status_code=403,
            detail="Export limit reached. Upgrade to Pro for more exports."
        )

    try:
        import pandas as pd
    except ImportError:
        raise HTTPException(status_code=500, detail="CSV export requires pandas")

    # Fetch data
    if data.entity_type == "drugs":
        df = await fetch_drugs_data(data.entity_ids, data.include_scores)
        filename = "phase4_drugs_export"
    elif data.entity_type == "targets":
        df = await fetch_targets_data(data.entity_ids)
        filename = "phase4_targets_export"
    elif data.entity_type == "trials":
        df = await fetch_trials_data(data.entity_ids)
        filename = "phase4_trials_export"
    elif data.entity_type == "scores":
        df = await fetch_scores_data(data.entity_ids)
        filename = "phase4_scores_export"
    else:
        raise HTTPException(status_code=400, detail=f"Invalid entity type: {data.entity_type}")

    if df.empty:
        raise HTTPException(status_code=404, detail="No data found for export")

    # Convert to CSV
    output = io.StringIO()
    df.to_csv(output, index=False)
    output.seek(0)

    # Increment export count
    await increment_export_count(user_id)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return StreamingResponse(
        io.BytesIO(output.getvalue().encode()),
        media_type="text/csv",
        headers={
            "Content-Disposition": f"attachment; filename={filename}_{timestamp}.csv"
        }
    )


# ============================================
# PowerPoint Deal Memo
# ============================================

@router.post("/deal-memo")
async def generate_deal_memo(
    data: DealMemoRequest,
    request: Request,
    user_id: str = Depends(require_auth)
):
    """Generate a PowerPoint deal memo for a drug."""
    # Check export access
    if not await check_export_access(user_id):
        raise HTTPException(
            status_code=403,
            detail="Export limit reached. Upgrade to Pro for more exports."
        )

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint export requires python-pptx. Run: pip install python-pptx"
        )

    # Fetch drug data
    drug = supabase.table("epi_drugs").select("*").eq("id", data.drug_id).execute()
    if not drug.data:
        raise HTTPException(status_code=404, detail="Drug not found")
    drug = drug.data[0]

    # Fetch scores
    scores = supabase.table("epi_scores").select(
        "*, epi_indications(name)"
    ).eq("drug_id", data.drug_id).execute()

    # Fetch targets
    targets = supabase.table("epi_drug_targets").select(
        "*, epi_targets(symbol, name, family)"
    ).eq("drug_id", data.drug_id).execute()

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Dark background
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(0, 0, 0)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = drug["name"]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    target_names = ", ".join([t["epi_targets"]["symbol"] for t in (targets.data or [])[:3]])
    subtitle_para.text = f"{target_names} Inhibitor | Phase {drug.get('max_phase', 'N/A')}"
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RgbColor(160, 160, 160)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = f"Phase4 Intelligence | Deal Memo | {datetime.now().strftime('%B %Y')}"
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = RgbColor(100, 100, 100)
    footer_para.alignment = PP_ALIGN.CENTER

    # Slide 2: Overview
    slide = prs.slides.add_slide(slide_layout)
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RgbColor(0, 0, 0)
    background.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Asset Overview"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RgbColor(255, 255, 255)

    # Key metrics
    metrics = [
        ("Drug Name", drug["name"]),
        ("Drug Type", drug.get("drug_type", "Small molecule")),
        ("Max Phase", str(drug.get("max_phase", "N/A"))),
        ("FDA Approved", "Yes" if drug.get("fda_approved") else "No"),
        ("ChEMBL ID", drug.get("chembl_id", "N/A")),
        ("Targets", target_names or "N/A"),
    ]

    y_pos = 1.5
    for label, value in metrics:
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(10), Inches(0.4))
        frame = box.text_frame
        para = frame.paragraphs[0]
        para.text = f"{label}: {value}"
        para.font.size = Pt(18)
        para.font.color.rgb = RgbColor(200, 200, 200)
        y_pos += 0.5

    # Slide 3: Scores
    if scores.data:
        slide = prs.slides.add_slide(slide_layout)
        background = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = RgbColor(0, 0, 0)
        background.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "TotalScore Breakdown"
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RgbColor(255, 255, 255)

        y_pos = 1.5
        for score in scores.data[:5]:
            indication_name = score.get("epi_indications", {}).get("name", "Unknown")
            total = score.get("total_score", 0) or 0
            bio = score.get("bio_score", 0) or 0
            chem = score.get("chem_score", 0) or 0
            tract = score.get("tractability_score", 0) or 0

            box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(10), Inches(0.8))
            frame = box.text_frame
            para = frame.paragraphs[0]
            para.text = f"{indication_name}"
            para.font.size = Pt(20)
            para.font.bold = True
            para.font.color.rgb = RgbColor(255, 255, 255)

            para2 = frame.add_paragraph()
            para2.text = f"Total: {total:.1f} | Bio: {bio:.1f} | Chem: {chem:.1f} | Tract: {tract:.1f}"
            para2.font.size = Pt(14)
            para2.font.color.rgb = RgbColor(150, 150, 150)

            y_pos += 1.0

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # Increment export count
    await increment_export_count(user_id)

    # Return file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_name = drug["name"].replace(" ", "_").replace("/", "-")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=Phase4_DealMemo_{safe_name}_{timestamp}.pptx"
        }
    )


# ============================================
# Competitive Landscape PPTX Export
# ============================================

@router.post("/landscape")
async def generate_landscape_export(
    data: LandscapeExportRequest,
    request: Request,
    user_id: str = Depends(require_auth)
):
    """
    Generate a PowerPoint competitive landscape deck.

    Supports multiple export types:
    - target: All drugs targeting a specific protein (e.g., HDAC1 landscape)
    - indication: All drugs in a specific indication (e.g., AML landscape)
    - company: Company portfolio with all drugs and editing assets
    - pipeline: Custom selection of drugs for comparison
    """
    # Check export access
    if not await check_export_access(user_id):
        raise HTTPException(
            status_code=403,
            detail="Export limit reached. Upgrade to Team or Enterprise for more exports."
        )

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint export requires python-pptx. Run: pip install python-pptx"
        )

    # Fetch data based on export type
    if data.export_type == "target":
        if not data.target_id:
            raise HTTPException(status_code=400, detail="target_id required for target landscape")
        landscape_data = await fetch_target_landscape(data.target_id, data.include_scores, data.include_trials)
        title = f"{landscape_data['target']['symbol']} Competitive Landscape"
        subtitle = f"{landscape_data['target']['family']} | {len(landscape_data['drugs'])} Pipeline Assets"
        filename_prefix = f"Target_Landscape_{landscape_data['target']['symbol']}"

    elif data.export_type == "indication":
        if not data.indication_id:
            raise HTTPException(status_code=400, detail="indication_id required for indication landscape")
        landscape_data = await fetch_indication_landscape(data.indication_id, data.include_scores, data.include_trials)
        title = f"{landscape_data['indication']['name']} Competitive Landscape"
        subtitle = f"Epigenetic Drug Pipeline | {len(landscape_data['drugs'])} Assets"
        filename_prefix = f"Indication_Landscape_{landscape_data['indication']['name'].replace(' ', '_')}"

    elif data.export_type == "company":
        if not data.company_id:
            raise HTTPException(status_code=400, detail="company_id required for company portfolio")
        landscape_data = await fetch_company_landscape(data.company_id, data.include_scores, data.include_trials)
        title = f"{landscape_data['company']['name']} Pipeline Portfolio"
        ticker = landscape_data['company'].get('ticker') or "Private"
        subtitle = f"{ticker} | {len(landscape_data['drugs'])} Drugs | {len(landscape_data.get('editing_assets', []))} Editing Programs"
        filename_prefix = f"Company_Portfolio_{landscape_data['company']['name'].replace(' ', '_')}"

    elif data.export_type == "pipeline":
        if not data.drug_ids or len(data.drug_ids) == 0:
            raise HTTPException(status_code=400, detail="drug_ids required for pipeline comparison")
        landscape_data = await fetch_pipeline_comparison(data.drug_ids, data.include_scores, data.include_trials)
        title = "Pipeline Comparison Analysis"
        subtitle = f"{len(landscape_data['drugs'])} Selected Assets"
        filename_prefix = "Pipeline_Comparison"
    else:
        raise HTTPException(status_code=400, detail=f"Invalid export_type: {data.export_type}")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    BG_COLOR = RgbColor(0, 0, 0)
    TEXT_PRIMARY = RgbColor(255, 255, 255)
    TEXT_SECONDARY = RgbColor(160, 160, 160)
    TEXT_MUTED = RgbColor(100, 100, 100)
    ACCENT_BLUE = RgbColor(96, 165, 250)  # blue-400
    SCORE_HIGH = RgbColor(34, 197, 94)    # green-500
    SCORE_MED = RgbColor(234, 179, 8)     # yellow-500
    SCORE_LOW = RgbColor(148, 163, 184)   # slate-400

    def add_dark_background(slide):
        """Add dark background to slide."""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = BG_COLOR
        background.line.fill.background()
        # Send to back
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def get_score_color(score):
        """Get color based on score value."""
        if score is None:
            return TEXT_MUTED
        if score >= 60:
            return SCORE_HIGH
        if score >= 40:
            return SCORE_MED
        return SCORE_LOW

    # ===== Slide 1: Title Slide =====
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Phase4 Intelligence | {datetime.now().strftime('%B %Y')}"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # ===== Slide 2: Executive Summary =====
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY

    # Summary stats
    drugs = landscape_data.get('drugs', [])
    scores = [d.get('total_score') for d in drugs if d.get('total_score')]
    avg_score = sum(scores) / len(scores) if scores else 0

    phase_counts = {}
    for d in drugs:
        phase = d.get('max_phase') or d.get('phase') or 0
        phase_counts[phase] = phase_counts.get(phase, 0) + 1

    approved_count = sum(1 for d in drugs if d.get('fda_approved'))

    stats = [
        ("Total Assets", str(len(drugs))),
        ("Avg TotalScore", f"{avg_score:.1f}" if avg_score else "N/A"),
        ("FDA Approved", str(approved_count)),
        ("Phase 3", str(phase_counts.get(3, 0))),
        ("Phase 2", str(phase_counts.get(2, 0))),
        ("Phase 1", str(phase_counts.get(1, 0))),
    ]

    # Stats grid
    x_start = 0.75
    y_start = 1.5
    box_width = 2.0
    box_height = 1.2

    for i, (label, value) in enumerate(stats):
        col = i % 3
        row = i // 3
        x = x_start + col * (box_width + 0.3)
        y = y_start + row * (box_height + 0.3)

        # Stat box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(26, 26, 26)
        box.line.color.rgb = RgbColor(51, 51, 51)

        # Value
        val_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.2), Inches(box_width), Inches(0.6))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(Inches(x), Inches(y + 0.75), Inches(box_width), Inches(0.4))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT_MUTED
        p.alignment = PP_ALIGN.CENTER

    # ===== Slide 3: Pipeline Table =====
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Pipeline Assets"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY

    # Table data
    table_data = []
    for d in sorted(drugs, key=lambda x: x.get('total_score') or 0, reverse=True)[:15]:
        target_name = d.get('target_symbol') or d.get('target_name') or "-"
        table_data.append([
            d.get('name', 'Unknown'),
            target_name[:12],
            d.get('drug_type', '-')[:15] if d.get('drug_type') else '-',
            str(d.get('max_phase') or d.get('phase') or 0),
            f"{d.get('total_score', 0):.1f}" if d.get('total_score') else "-",
        ])

    if table_data:
        # Create table
        rows = min(len(table_data) + 1, 16)  # Header + data rows
        cols = 5
        table_width = Inches(12)
        table_height = Inches(5.5)

        table = slide.shapes.add_table(rows, cols, Inches(0.65), Inches(1.0), table_width, table_height).table

        # Set column widths
        table.columns[0].width = Inches(3.5)  # Drug
        table.columns[1].width = Inches(2.0)  # Target
        table.columns[2].width = Inches(2.5)  # Type
        table.columns[3].width = Inches(1.5)  # Phase
        table.columns[4].width = Inches(2.5)  # Score

        # Header row
        headers = ["Drug", "Target", "Type", "Phase", "TotalScore"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(26, 26, 26)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = TEXT_SECONDARY
            p.alignment = PP_ALIGN.CENTER

        # Data rows
        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = cell_text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(17, 17, 17) if row_idx % 2 == 0 else RgbColor(10, 10, 10)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = TEXT_PRIMARY if col_idx == 0 else TEXT_SECONDARY
                p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

                # Color score column
                if col_idx == 4 and cell_text != "-":
                    try:
                        score = float(cell_text)
                        p.font.color.rgb = get_score_color(score)
                        p.font.bold = True
                    except ValueError:
                        pass

    # ===== Slide 4: Score Distribution (if scores included) =====
    if data.include_scores and scores:
        slide = prs.slides.add_slide(slide_layout)
        add_dark_background(slide)

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Score Distribution"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY

        # Score breakdown by tier
        high_scores = [s for s in scores if s >= 60]
        med_scores = [s for s in scores if 40 <= s < 60]
        low_scores = [s for s in scores if s < 40]

        tiers = [
            ("High (≥60)", len(high_scores), SCORE_HIGH, high_scores[:3]),
            ("Medium (40-59)", len(med_scores), SCORE_MED, med_scores[:3]),
            ("Low (<40)", len(low_scores), SCORE_LOW, low_scores[:3]),
        ]

        y_pos = 1.2
        for tier_name, count, color, examples in tiers:
            # Tier label
            box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(4), Inches(0.5))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{tier_name}: {count} drugs"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = color

            # Bar visualization
            bar_width = min(count / len(scores) * 8, 8) if scores else 0
            if bar_width > 0:
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(5), Inches(y_pos + 0.1), Inches(bar_width), Inches(0.4)
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = color
                bar.line.fill.background()

            y_pos += 1.2

        # Top performers
        top_drugs = sorted(drugs, key=lambda x: x.get('total_score') or 0, reverse=True)[:5]
        if top_drugs:
            title_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(10), Inches(0.5))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = "Top Performing Assets"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = TEXT_PRIMARY

            y_pos = 5.4
            for d in top_drugs:
                box = slide.shapes.add_textbox(Inches(1.5), Inches(y_pos), Inches(10), Inches(0.4))
                tf = box.text_frame
                p = tf.paragraphs[0]
                score = d.get('total_score', 0)
                p.text = f"• {d.get('name', 'Unknown')} — {score:.1f}"
                p.font.size = Pt(14)
                p.font.color.rgb = get_score_color(score)
                y_pos += 0.35

    # ===== Slide 5: Key Takeaways =====
    slide = prs.slides.add_slide(slide_layout)
    add_dark_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Takeaways"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY

    takeaways = []
    if len(drugs) > 0:
        takeaways.append(f"Pipeline contains {len(drugs)} assets across multiple development phases")
    if approved_count > 0:
        takeaways.append(f"{approved_count} FDA-approved drugs provide revenue/validation")
    if phase_counts.get(3, 0) > 0:
        takeaways.append(f"{phase_counts.get(3, 0)} Phase 3 programs represent near-term catalysts")
    if avg_score >= 50:
        takeaways.append(f"Strong average TotalScore ({avg_score:.1f}) indicates compelling biology + chemistry")
    elif avg_score > 0:
        takeaways.append(f"Moderate average TotalScore ({avg_score:.1f}) suggests room for optimization")

    # Add context-specific takeaway
    if data.export_type == "target":
        target = landscape_data.get('target', {})
        if target.get('io_exhaustion_axis'):
            takeaways.append(f"{target.get('symbol')} is implicated in IO exhaustion — combo opportunity with checkpoint inhibitors")

    y_pos = 1.2
    for i, takeaway in enumerate(takeaways[:6]):
        box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11), Inches(0.7))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {takeaway}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SECONDARY
        y_pos += 0.9

    # Footer with disclaimer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Source: Phase4 Intelligence | Data as of " + datetime.now().strftime('%Y-%m-%d')
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    # Increment export count
    await increment_export_count(user_id)

    # Return file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename=Phase4_{filename_prefix}_{timestamp}.pptx"
        }
    )


# ============================================
# Landscape Data Fetching
# ============================================

async def fetch_target_landscape(target_id: str, include_scores: bool, include_trials: bool):
    """Fetch all drugs targeting a specific protein."""
    # Get target info
    target = supabase.table("epi_targets").select("*").eq("id", target_id).single().execute()
    if not target.data:
        raise HTTPException(status_code=404, detail="Target not found")

    # Get drugs via drug_targets
    drug_links = supabase.table("epi_drug_targets").select(
        "drug_id, mechanism_of_action, epi_drugs(*)"
    ).eq("target_id", target_id).execute()

    drugs = []
    for link in (drug_links.data or []):
        drug = link.get("epi_drugs", {})
        if drug:
            drug_data = {
                **drug,
                "mechanism": link.get("mechanism_of_action"),
                "target_symbol": target.data.get("symbol"),
            }

            # Get scores
            if include_scores:
                scores = supabase.table("epi_scores").select("*").eq("drug_id", drug["id"]).execute()
                if scores.data:
                    best_score = max(scores.data, key=lambda s: s.get("total_score") or 0)
                    drug_data.update({
                        "total_score": best_score.get("total_score"),
                        "bio_score": best_score.get("bio_score"),
                        "chem_score": best_score.get("chem_score"),
                        "tractability_score": best_score.get("tractability_score"),
                    })

            drugs.append(drug_data)

    return {
        "target": target.data,
        "drugs": drugs,
    }


async def fetch_indication_landscape(indication_id: str, include_scores: bool, include_trials: bool):
    """Fetch all drugs in a specific indication."""
    # Get indication info
    indication = supabase.table("epi_indications").select("*").eq("id", indication_id).single().execute()
    if not indication.data:
        raise HTTPException(status_code=404, detail="Indication not found")

    # Get drugs via drug_indications
    drug_links = supabase.table("epi_drug_indications").select(
        "drug_id, approval_status, epi_drugs(*)"
    ).eq("indication_id", indication_id).execute()

    drugs = []
    for link in (drug_links.data or []):
        drug = link.get("epi_drugs", {})
        if drug:
            drug_data = {
                **drug,
                "approval_status": link.get("approval_status"),
            }

            # Get target info
            targets = supabase.table("epi_drug_targets").select(
                "epi_targets(symbol)"
            ).eq("drug_id", drug["id"]).limit(1).execute()
            if targets.data and targets.data[0].get("epi_targets"):
                drug_data["target_symbol"] = targets.data[0]["epi_targets"]["symbol"]

            # Get scores for this indication
            if include_scores:
                scores = supabase.table("epi_scores").select("*").eq(
                    "drug_id", drug["id"]
                ).eq("indication_id", indication_id).execute()
                if scores.data:
                    score = scores.data[0]
                    drug_data.update({
                        "total_score": score.get("total_score"),
                        "bio_score": score.get("bio_score"),
                        "chem_score": score.get("chem_score"),
                        "tractability_score": score.get("tractability_score"),
                    })

            drugs.append(drug_data)

    return {
        "indication": indication.data,
        "drugs": drugs,
    }


async def fetch_company_landscape(company_id: str, include_scores: bool, include_trials: bool):
    """Fetch company portfolio with drugs and editing assets."""
    # Get company info
    company = supabase.table("epi_companies").select("*").eq("id", company_id).single().execute()
    if not company.data:
        raise HTTPException(status_code=404, detail="Company not found")

    # Get drugs
    drug_links = supabase.table("epi_drug_companies").select(
        "drug_id, role, is_primary, epi_drugs(*)"
    ).eq("company_id", company_id).execute()

    drugs = []
    for link in (drug_links.data or []):
        drug = link.get("epi_drugs", {})
        if drug:
            drug_data = {
                **drug,
                "role": link.get("role"),
                "is_primary": link.get("is_primary"),
            }

            # Get target info
            targets = supabase.table("epi_drug_targets").select(
                "epi_targets(symbol)"
            ).eq("drug_id", drug["id"]).limit(1).execute()
            if targets.data and targets.data[0].get("epi_targets"):
                drug_data["target_symbol"] = targets.data[0]["epi_targets"]["symbol"]

            # Get scores
            if include_scores:
                scores = supabase.table("epi_scores").select("*").eq("drug_id", drug["id"]).execute()
                if scores.data:
                    best_score = max(scores.data, key=lambda s: s.get("total_score") or 0)
                    drug_data.update({
                        "total_score": best_score.get("total_score"),
                        "bio_score": best_score.get("bio_score"),
                        "chem_score": best_score.get("chem_score"),
                        "tractability_score": best_score.get("tractability_score"),
                    })

            drugs.append(drug_data)

    # Get editing assets
    editing_links = supabase.table("epi_editing_asset_companies").select(
        "editing_asset_id, role, is_primary, epi_editing_assets(*)"
    ).eq("company_id", company_id).execute()

    editing_assets = []
    for link in (editing_links.data or []):
        asset = link.get("epi_editing_assets", {})
        if asset:
            asset_data = {
                **asset,
                "role": link.get("role"),
                "is_primary": link.get("is_primary"),
            }
            editing_assets.append(asset_data)

    return {
        "company": company.data,
        "drugs": drugs,
        "editing_assets": editing_assets,
    }


async def fetch_pipeline_comparison(drug_ids: List[str], include_scores: bool, include_trials: bool):
    """Fetch specific drugs for comparison."""
    drugs = []

    for drug_id in drug_ids:
        drug_result = supabase.table("epi_drugs").select("*").eq("id", drug_id).single().execute()
        if drug_result.data:
            drug_data = drug_result.data

            # Get target info
            targets = supabase.table("epi_drug_targets").select(
                "epi_targets(symbol)"
            ).eq("drug_id", drug_id).limit(1).execute()
            if targets.data and targets.data[0].get("epi_targets"):
                drug_data["target_symbol"] = targets.data[0]["epi_targets"]["symbol"]

            # Get scores
            if include_scores:
                scores = supabase.table("epi_scores").select("*").eq("drug_id", drug_id).execute()
                if scores.data:
                    best_score = max(scores.data, key=lambda s: s.get("total_score") or 0)
                    drug_data.update({
                        "total_score": best_score.get("total_score"),
                        "bio_score": best_score.get("bio_score"),
                        "chem_score": best_score.get("chem_score"),
                        "tractability_score": best_score.get("tractability_score"),
                    })

            drugs.append(drug_data)

    return {
        "drugs": drugs,
    }


# ============================================
# Data Fetching Helpers
# ============================================

async def fetch_drugs_data(entity_ids: Optional[List[str]], include_scores: bool):
    """Fetch drugs data for export."""
    import pandas as pd

    query = supabase.table("epi_drugs").select(
        "id, name, chembl_id, drug_type, max_phase, fda_approved, modality"
    )

    if entity_ids:
        query = query.in_("id", entity_ids)

    result = query.order("name").execute()
    drugs = result.data or []

    if include_scores and drugs:
        # Get scores for each drug
        drug_ids = [d["id"] for d in drugs]
        scores = supabase.table("epi_scores").select(
            "drug_id, total_score, bio_score, chem_score, tractability_score"
        ).in_("drug_id", drug_ids).execute()

        score_map = {}
        for s in (scores.data or []):
            if s["drug_id"] not in score_map or (s.get("total_score") or 0) > (score_map[s["drug_id"]].get("total_score") or 0):
                score_map[s["drug_id"]] = s

        for drug in drugs:
            score = score_map.get(drug["id"], {})
            drug["total_score"] = score.get("total_score")
            drug["bio_score"] = score.get("bio_score")
            drug["chem_score"] = score.get("chem_score")
            drug["tractability_score"] = score.get("tractability_score")

    return pd.DataFrame(drugs)


async def fetch_targets_data(entity_ids: Optional[List[str]]):
    """Fetch targets data for export."""
    import pandas as pd

    query = supabase.table("epi_targets").select(
        "id, symbol, name, family, target_class, ensembl_gene_id, uniprot_id, io_exhaustion_axis, io_combo_priority"
    )

    if entity_ids:
        query = query.in_("id", entity_ids)

    result = query.order("symbol").execute()
    return pd.DataFrame(result.data or [])


async def fetch_trials_data(entity_ids: Optional[List[str]]):
    """Fetch trials data for export."""
    import pandas as pd

    query = supabase.table("ci_trial_calendar").select(
        "nct_id, trial_title, drug_name, phase, status, primary_completion_date, enrollment, lead_sponsor"
    )

    if entity_ids:
        query = query.in_("nct_id", entity_ids)

    result = query.order("primary_completion_date").execute()
    return pd.DataFrame(result.data or [])


async def fetch_scores_data(entity_ids: Optional[List[str]]):
    """Fetch scores data for export."""
    import pandas as pd

    query = supabase.table("epi_scores").select(
        "*, epi_drugs(name), epi_indications(name)"
    )

    if entity_ids:
        query = query.in_("drug_id", entity_ids)

    result = query.order("total_score", desc=True).execute()

    # Flatten the data
    scores = []
    for s in (result.data or []):
        scores.append({
            "drug_name": s.get("epi_drugs", {}).get("name"),
            "indication": s.get("epi_indications", {}).get("name"),
            "total_score": s.get("total_score"),
            "bio_score": s.get("bio_score"),
            "chem_score": s.get("chem_score"),
            "tractability_score": s.get("tractability_score"),
        })

    return pd.DataFrame(scores)


async def fetch_watchlist_data(user_id: str):
    """Fetch user's watchlist for export."""
    import pandas as pd

    result = supabase.table("ci_watchlist").select(
        "entity_type, entity_name, notes, created_at"
    ).eq("user_id", user_id).execute()

    return pd.DataFrame(result.data or [])
